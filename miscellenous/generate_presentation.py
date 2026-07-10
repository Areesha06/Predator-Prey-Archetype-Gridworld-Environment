#!/usr/bin/env python3
"""Generate presentation.pptx for the Predator-Prey MARL Gridworld project."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
TEAL = RGBColor(0x0D, 0x6E, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xED, 0xED, 0xED)
MGRAY = RGBColor(0x88, 0x88, 0x88)
DGRAY = RGBColor(0x2D, 0x2D, 0x2D)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x1A, 0x7A, 0x40)
BLUE = RGBColor(0x21, 0x6C, 0xB4)
PURPLE = RGBColor(0x6C, 0x3B, 0x9E)
AMBER = RGBColor(0xF1, 0x9C, 0x12)
CODEBG = RGBColor(0xED, 0xED, 0xED)

# Layer colours for 4-layer diagram
L1 = RGBColor(0xC0, 0x39, 0x2B)
L2 = RGBColor(0xE6, 0x7E, 0x22)
L3 = RGBColor(0x1A, 0x7A, 0x40)
L4 = RGBColor(0x21, 0x6C, 0xB4)

# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Primitive helpers ─────────────────────────────────────────────────────────


def S():
    return prs.slides.add_slide(BLANK)


def bg(sl, c):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = c


def rect(sl, l, t, w, h, fill=None, border=None, bw=0.75):
    shp = sl.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp


def txt(
    sl,
    text,
    l,
    t,
    w,
    h,
    size=16,
    bold=False,
    italic=False,
    color=None,
    align=PP_ALIGN.LEFT,
    font="Calibri",
):
    color = color or DGRAY
    txb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return txb


def mltext(sl, lines, l, t, w, h, def_size=15, def_color=None, font="Calibri"):
    """lines = list of str | (str, bold, size, color)"""
    def_color = def_color or DGRAY
    txb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            s, b, sz, c = item, False, def_size, def_color
        else:
            s = item[0]
            b = item[1] if len(item) > 1 else False
            sz = item[2] if len(item) > 2 else def_size
            c = item[3] if len(item) > 3 else def_color
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = s
        r.font.size = Pt(sz)
        r.font.bold = b
        r.font.color.rgb = c
        r.font.name = font
    return txb


def code(sl, src, l, t, w, h, size=11):
    rect(sl, l, t, w, h, fill=CODEBG)
    txb = sl.shapes.add_textbox(
        Inches(l + 0.12), Inches(t + 0.08), Inches(w - 0.24), Inches(h - 0.16)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in src.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.name = "Courier New"
        r.font.color.rgb = DGRAY


def hdr(sl, title, hc=NAVY, tc=WHITE, h=0.85):
    rect(sl, 0, 0, 13.33, h, fill=hc)
    txt(sl, title, 0.35, 0.10, 12.6, h - 0.18, size=24, bold=True, color=tc)


def tbl(
    sl, data, l, t, w, h, hfill=NAVY, htxt=WHITE, col_w=None, fs=13, row_fills=None
):
    rows, cols = len(data), len(data[0])
    t_ = sl.shapes.add_table(
        rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)
    ).table
    if col_w:
        emu = Inches(w)
        for i, f in enumerate(col_w):
            t_.columns[i].width = int(emu * f)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = t_.cell(ri, ci)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(fs)
                    run.font.name = "Calibri"
                    run.font.bold = ri == 0
                    run.font.color.rgb = htxt if ri == 0 else DGRAY
            f = cell.fill
            if ri == 0:
                f.solid()
                f.fore_color.rgb = hfill
            elif row_fills and ri - 1 < len(row_fills) and row_fills[ri - 1]:
                f.solid()
                f.fore_color.rgb = row_fills[ri - 1]


# ── Compound helpers ──────────────────────────────────────────────────────────


def section(title, sub=""):
    sl = S()
    bg(sl, NAVY)
    rect(sl, 1.0, 3.35, 11.33, 0.07, fill=ORANGE)
    txt(
        sl,
        title,
        0.5,
        1.9,
        12.33,
        1.3,
        size=38,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    if sub:
        txt(
            sl,
            sub,
            0.5,
            3.55,
            12.33,
            0.8,
            size=20,
            color=RGBColor(0xB0, 0xC4, 0xDE),
            align=PP_ALIGN.CENTER,
        )
    return sl


def cslide(title, hc=NAVY):
    sl = S()
    bg(sl, WHITE)
    hdr(sl, title, hc=hc)
    return sl


def bullets(sl, items, l, t, w, h, size=16, color=None, font="Calibri"):
    color = color or DGRAY
    txb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        txt_s, lvl = (item, 0) if isinstance(item, str) else (item[0], item[1])
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        sym = "•  " if lvl == 0 else "◦  "
        ind = "    " * lvl
        r = p.add_run()
        r.text = ind + sym + txt_s
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = color
        r.font.name = font


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl, NAVY)
rect(sl, 0, 0, 13.33, 0.32, fill=ORANGE)
txt(
    sl,
    "Predator-Prey Gridworld",
    0.6,
    1.1,
    12.13,
    1.4,
    size=44,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)
txt(
    sl,
    "A Modular Multi-Agent Reinforcement Learning Testbed",
    0.6,
    2.75,
    12.13,
    0.8,
    size=24,
    color=RGBColor(0xB0, 0xC4, 0xDE),
    align=PP_ALIGN.CENTER,
)
rect(sl, 1.0, 3.7, 11.33, 0.05, fill=ORANGE)
txt(
    sl,
    "Student Orientation  •  Week 1",
    0.6,
    3.95,
    12.13,
    0.5,
    size=18,
    color=RGBColor(0x90, 0xA8, 0xC8),
    align=PP_ALIGN.CENTER,
)
txt(
    sl,
    "Predator-Prey Archetype Gridworld Environment",
    0.6,
    4.55,
    12.13,
    0.5,
    size=15,
    color=MGRAY,
    align=PP_ALIGN.CENTER,
)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 1 — Why This Exists
# ══════════════════════════════════════════════════════════════════════════════
section("Section 1: Why This Exists", "The problem we solve and our mission")

# Slide 2 — The Problem
sl = cslide("The Problem with Most MARL Environments")
txt(
    sl,
    "Most research environments tangle concerns that should be independent:",
    0.5,
    1.0,
    12.33,
    0.45,
    size=16,
    color=DGRAY,
)
bullets(
    sl,
    [
        "Environment dynamics and learning algorithm live in the same class",
        "Opaque reward computation — you can't trace why a policy improves",
        "Global random state makes experiments non-reproducible across machines",
        "Changing one variable (e.g. reward shape) silently affects others",
        "Hard to run controlled ablations — observation vs reward vs algorithm",
    ],
    0.5,
    1.55,
    12.33,
    3.2,
    size=17,
)
rect(sl, 0.5, 5.0, 12.33, 1.6, fill=RGBColor(0xE8, 0xF4, 0xFD), border=BLUE)
txt(
    sl,
    '"We need an environment where each layer can be studied independently,\n'
    'with full reproducibility and zero surprise interactions between components."',
    0.8,
    5.15,
    11.8,
    1.3,
    size=16,
    italic=True,
    color=NAVY,
)

# Slide 3 — Is / Is Not
sl = cslide("What This Project Is — and Is Not")
rect(sl, 0.4, 1.0, 5.85, 5.8, fill=RGBColor(0xF0, 0xFB, 0xF4), border=GREEN)
txt(sl, "IS", 0.65, 1.05, 5.3, 0.5, size=20, bold=True, color=GREEN)
y = 1.65
for s in [
    "A controlled lab for teaching and researching RL/MARL",
    "Config-driven — one YAML change = one experimental variable",
    "Fully deterministic — same seed, same trajectory, always",
    "Plugin-based — swap observations, rewards, and algorithms",
    "Traceable — every transition can be verified by hand",
]:
    txt(sl, "✓  " + s, 0.65, y, 5.3, 0.55, size=14, color=DGRAY)
    y += 0.7

rect(sl, 6.9, 1.0, 5.85, 5.8, fill=RGBColor(0xFE, 0xF5, 0xF5), border=RED)
txt(sl, "IS NOT", 7.15, 1.05, 5.3, 0.5, size=20, bold=True, color=RED)
y = 1.65
for s in [
    "A high-performance training platform",
    "A continuous-action or visual environment",
    "A replacement for PettingZoo or SMAC",
    "Optimised for GPU-scale deep RL",
    "A ready-made publication benchmark",
]:
    txt(sl, "✗  " + s, 7.15, y, 5.3, 0.55, size=14, color=DGRAY)
    y += 0.7

txt(
    sl,
    "← designed for this",
    3.5,
    6.6,
    3.0,
    0.4,
    size=13,
    color=GREEN,
    align=PP_ALIGN.RIGHT,
)
txt(sl, "out of scope →", 7.15, 6.6, 3.0, 0.4, size=13, color=RED)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
section("Section 2: Architecture", "Four layers, one golden rule")

# Slide 4 — Four-Layer Diagram
sl = cslide("The Four-Layer Architecture")
txt(
    sl,
    "Every component lives in exactly one layer. Dependencies flow downward only.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    color=MGRAY,
)
layers = [
    (
        L4,
        "Layer 4  —  Scripts & Orchestration",
        "run_from_config.py  •  evaluate.py  •  sweep.py  •  render.py",
        "Drives training runs from YAML config",
    ),
    (
        L3,
        "Layer 3  —  Baselines (Learning Algorithms)",
        "IQL  •  CQL  •  MixedTrainer  •  (future: DQN, PPO, Nash Q)",
        "Selects actions, updates Q-tables",
    ),
    (
        L2,
        "Layer 2  —  Plugins (Observations + Rewards)",
        "5 observation builders  •  3 reward functions",
        "Transforms env state into percepts and signals",
    ),
    (
        L1,
        "Layer 1  —  Core Environment   ★ IMMUTABLE ★",
        "core/gridworld.py  •  core/agent.py",
        "Grid physics, movement, capture — never modified",
    ),
]
y = 1.45
for col, title, sub, note in layers:
    rect(sl, 0.4, y, 12.5, 1.15, fill=col)
    txt(sl, title, 0.65, y + 0.08, 8.5, 0.45, size=17, bold=True, color=WHITE)
    txt(sl, sub, 0.65, y + 0.55, 8.5, 0.38, size=13, color=RGBColor(0xDD, 0xDD, 0xDD))
    lighter = RGBColor(
        min(255, col[0] + 55), min(255, col[1] + 55), min(255, col[2] + 55)
    )
    rect(sl, 9.0, y + 0.08, 3.65, 0.97, fill=lighter)
    txt(sl, note, 9.1, y + 0.25, 3.45, 0.7, size=12, color=WHITE)
    y += 1.22

# Slide 5 — Design Invariants
sl = cslide("Five Core Design Invariants")
txt(
    sl,
    "These rules protect the scientific validity of every experiment.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    color=MGRAY,
)
invariants = [
    (
        "Immutable Core",
        "src/core/ is never modified by contributors. Physics must not change between experiments.",
    ),
    (
        "Seeded Randomness",
        "env.reset(seed=n) with the same n always produces the exact same trajectory.",
    ),
    (
        "Read-Only Plugins",
        "Observation builders and reward functions observe state — they never write it.",
    ),
    (
        "Config Drives All",
        "If you need behavioral change, write a plugin. If you think you need to edit core, that is a design bug.",
    ),
    (
        "Black-Box Algos",
        "Algorithms may only call env.reset() and env.step(). No peeking at env.agents or internal state.",
    ),
]
y = 1.45
for title, body in invariants:
    rect(sl, 0.4, y, 12.5, 0.9, fill=LGRAY, border=RGBColor(0xCC, 0xCC, 0xCC))
    txt(sl, title, 0.65, y + 0.1, 3.3, 0.5, size=15, bold=True, color=NAVY)
    txt(sl, body, 4.1, y + 0.12, 8.7, 0.72, size=14, color=DGRAY)
    y += 0.97

# Slide 6 — Plugin & Registry Pattern
sl = cslide("The Plugin & Registry Pattern")
txt(
    sl,
    "Plugins are discovered at runtime from YAML keys — no hardcoded imports in core code.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    color=MGRAY,
)
boxes = [
    (0.4, "YAML Key", "type: local_radius", BLUE),
    (3.0, "Registry", "get_observation_\nbuilder(key)", TEAL),
    (5.6, "Instantiation", "LocalRadius\nObs(**params)", GREEN),
    (8.2, "Injection", "env.observation_\nbuilder = fn", ORANGE),
    (10.8, "Called Each Step", "self.observation_\nbuilder(self)", RED),
]
for x, title, body, col in boxes:
    rect(sl, x, 1.55, 2.2, 1.55, fill=col)
    txt(
        sl,
        title,
        x + 0.1,
        1.6,
        2.0,
        0.42,
        size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    txt(
        sl,
        body,
        x + 0.1,
        2.06,
        2.0,
        1.0,
        size=11,
        color=WHITE,
        align=PP_ALIGN.CENTER,
        font="Courier New",
    )
for xi in [2.62, 5.22, 7.82, 10.42]:
    txt(
        sl,
        "→",
        xi,
        2.1,
        0.4,
        0.4,
        size=22,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )

rect(sl, 0.4, 3.4, 12.5, 3.4, fill=RGBColor(0xF5, 0xF8, 0xFF), border=BLUE)
txt(sl, "Why this matters", 0.65, 3.5, 4.0, 0.45, size=16, bold=True, color=NAVY)
bullets(
    sl,
    [
        "Adding a new observation builder: one new file + one registry line — zero core changes",
        "Algorithm author never imports the builder directly — YAML is the only coupling point",
        "Core code never changes — all prior experiments remain valid after adding a plugin",
        "Pattern is identical across all three extension points (obs / reward / algorithm)",
    ],
    0.65,
    4.05,
    12.0,
    2.5,
    size=15,
)

# Slide 7 — Extension Points Table
sl = cslide("Extension Points — The Three Contribution Pathways")
tbl(
    sl,
    [
        ["Layer", "Interface Method", "Where", "Registry File", "YAML Key"],
        [
            "Observation",
            "ObservationBuilder.build(env) → dict",
            "src/observations/",
            "observation_registry.py",
            "observation.type",
        ],
        [
            "Reward",
            "RewardFunction.compute(env) → dict",
            "src/rewards/",
            "reward_registry.py",
            "rewards[].name",
        ],
        [
            "Algorithm",
            "BaseAlgorithm.select_actions()\n.train()",
            "src/baselines/",
            "algorithm_registry.py",
            "experiment.algorithm.name",
        ],
    ],
    0.4,
    1.1,
    12.5,
    5.6,
    hfill=NAVY,
    col_w=[0.13, 0.28, 0.20, 0.22, 0.17],
    fs=14,
)
txt(
    sl,
    "All three follow the identical discover-by-string-key pattern. New contributors only ever touch one pathway.",
    0.5,
    6.85,
    12.33,
    0.42,
    size=13,
    italic=True,
    color=MGRAY,
)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 3 — Workflows
# ══════════════════════════════════════════════════════════════════════════════
section("Section 3: Workflows", "Init, step-loop, and training loop in detail")

# Slide 8 — Initialization Flow
sl = cslide("Initialization Flow — From YAML to Running Environment")
txt(
    sl,
    "What happens between  python run_from_config.py  and the first  env.step() ?",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
init_steps = [
    (BLUE, "① Load Configs", "load_all_configs()\nreads 5 YAML files"),
    (TEAL, "② Build Agents", "build_agents(cfg)\n→ list[Agent]"),
    (GREEN, "③ Create Env", "GridWorldEnv\n(agents, **params)"),
    (ORANGE, "④ Wire Observations", "observation_builder\nbound as callable"),
    (PURPLE, "⑤ Wire Rewards", "reward_fn closure\nbound to env"),
    (RED, "⑥ Start Training", "AlgorithmClass\n(env, cfg).train()"),
]
bw, bh = 2.0, 1.35
for i, (col, title, body) in enumerate(init_steps):
    row, col_i = divmod(i, 3)
    x = 0.35 + col_i * (12.7 / 3)
    y = 1.45 + row * (bh + 0.28)
    rect(sl, x, y, bw, bh, fill=col)
    txt(sl, title, x + 0.12, y + 0.08, bw - 0.24, 0.42, size=14, bold=True, color=WHITE)
    txt(
        sl,
        body,
        x + 0.12,
        y + 0.55,
        bw - 0.24,
        0.78,
        size=12,
        color=WHITE,
        font="Courier New",
    )

rect(sl, 0.4, 5.5, 12.5, 0.62, fill=NAVY)
txt(
    sl,
    "Result: a fully wired environment — observations, rewards, and algorithm all connected without the algorithm knowing how.",
    0.65,
    5.58,
    12.0,
    0.48,
    size=14,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)

# Slide 9 — Per-Step Pipeline
sl = cslide("Per-Step Pipeline — Inside  env.step(actions)")
txt(
    sl,
    "Called once per timestep. Every stage is deterministic given the same state and seed.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
steps9 = [
    "① Reset per-step trackers          ( _captured_this_step = [] )",
    "② Move all agents simultaneously   ( based on actions dict )",
    "③ Detect captures                  ( predator + prey on same cell → capture )",
    "④ Compute rewards                  ( reward_fn(self) → dict[name → float] )",
    "⑤ Check termination / truncation   ( all prey captured OR max_steps reached )",
    "⑥ Render frame                     ( only if render_mode='human' or 'rgb_array' )",
    "⑦ Build observations               ( observation_builder(self) → dict[name → dict] )",
    "⑧ Return result dict               ( obs, reward, terminated, truncated, info )",
]
y = 1.42
for i, s in enumerate(steps9):
    fill = LGRAY if i % 2 == 0 else WHITE
    rect(sl, 0.4, y, 12.5, 0.545, fill=fill, border=RGBColor(0xCC, 0xCC, 0xCC))
    txt(sl, s, 0.65, y + 0.1, 12.0, 0.38, size=14, color=DGRAY, font="Calibri")
    y += 0.56
rect(sl, 0.4, y + 0.02, 12.5, 0.5, fill=NAVY)
txt(
    sl,
    "⚠  Returns a dict — not the standard Gymnasium (obs, rew, term, trunc, info) tuple. Gymnasium wrapper is a Week-1 deliverable.",
    0.65,
    y + 0.1,
    12.0,
    0.38,
    size=13,
    color=AMBER,
)

# Slide 10 — Training Loop
sl = cslide("The Training Loop — IQL Episode Flow")
code(
    sl,
    """\
for episode in range(num_episodes):
    obs = env.reset()
    done = False

    while not done:
        actions = {}
        for agent_name, agent_obs in obs.items():
            state = _encode_state(agent_obs)
            if random() < epsilon:
                actions[agent_name] = randint(0, action_dim - 1)
            else:
                actions[agent_name] = argmax(q_tables[agent_name][state])

        result   = env.step(actions)

        for agent_name in agent_names:
            s  = _encode_state(obs[agent_name])
            a  = actions[agent_name]
            r  = result["reward"][agent_name]
            s_ = _encode_state(result["obs"][agent_name])
            td = r + gamma * max(q_tables[agent_name][s_].values()) \\
                   - q_tables[agent_name][s][a]
            q_tables[agent_name][s][a] += alpha * td

        obs  = result["obs"]
        done = result["terminated"] or result["truncated"]

    epsilon = max(min_epsilon, epsilon * epsilon_decay)
""",
    0.4,
    1.05,
    7.1,
    6.1,
    size=11,
)

annots = [
    (
        BLUE,
        "Epsilon-greedy",
        "Random with prob ε,\ngreedy otherwise.\nε decays per episode.",
    ),
    (
        TEAL,
        "Bellman Update",
        "Q(s,a) += α[r + γ·maxQ(s',·) − Q(s,a)]\nTD error drives convergence.",
    ),
    (
        GREEN,
        "defaultdict",
        "Q-tables auto-init new\nstates to zero. No\npre-allocation needed.",
    ),
    (
        ORANGE,
        "Epsilon Decay",
        "Exploration → exploitation.\nDecays after EVERY episode,\nnot every step.",
    ),
]
y = 1.1
for col, title, body in annots:
    rect(sl, 7.7, y, 5.3, 1.4, fill=col)
    txt(sl, title, 7.85, y + 0.07, 4.95, 0.42, size=14, bold=True, color=WHITE)
    txt(sl, body, 7.85, y + 0.55, 4.95, 0.82, size=12, color=WHITE)
    y += 1.5


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 4 — File Structure
# ══════════════════════════════════════════════════════════════════════════════
section("Section 4: File Structure", "Where everything lives and why")

# Slide 11 — Directory Tree
sl = cslide("Annotated Directory Structure")
txt(
    sl,
    "One directory = one responsibility. If something is hard to place, that is a design signal.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    color=MGRAY,
)
code(
    sl,
    """\
Predator-Prey-Archetype-Gridworld-Environment/
├─ src/
│  ├─ multi_agent_package/
│  │  ├─ core/           <- immutable physics (gridworld.py, agent.py)
│  │  ├─ observations/   <- 5 builders: LocalOnly, LocalRadius, Default, Absolute, Relative
│  │  ├─ rewards/        <- 3 functions: BaseReward, PredatorDistance, Survival
│  │  ├─ registry/       <- 3 registries: observation, reward (algorithm is in baselines)
│  │  └─ scripts/        <- run_from_config.py, evaluate.py, sweep.py, render.py
│  └─ baselines/
│     ├─ IQL/  CQL/  MIXED/   <- one folder per algorithm
│     ├─ base.py              <- BaseAlgorithm abstract class
│     └─ registry/            <- algorithm registry
├─ configs/              <- env.yaml  agents.yaml  observations.yaml  rewards.yaml  experiment*.yaml
├─ tests/                <- 189 pytest tests across 9 files (all architecture layers)
├─ notebooks/            <- 01_iql_demo.ipynb  02_cql_demo.ipynb
├─ wiki/                 <- 38 documentation files (SDD, ADRs, guides, concepts)
└─ .github/workflows/    <- ci.yaml (in progress)
""",
    0.4,
    1.42,
    8.6,
    5.55,
    size=12,
)

rect(sl, 9.3, 1.42, 3.8, 5.55, fill=RGBColor(0xF0, 0xF8, 0xFF), border=BLUE)
txt(sl, "Where do I...?", 9.5, 1.48, 3.4, 0.42, size=14, bold=True, color=NAVY)
guide = [
    ("Change how agents see", "observations/"),
    ("Change what agents want", "rewards/"),
    ("Add new algorithm", "baselines/"),
    ("Tune experiment params", "configs/"),
    ("Understand the system", "wiki/"),
    ("Run a demo interactively", "notebooks/"),
    ("Change movement rules", "talk to maintainer"),
]
y = 2.0
for q, a in guide:
    txt(sl, q, 9.5, y, 3.4, 0.3, size=12, color=DGRAY)
    txt(
        sl,
        "→  " + a,
        9.5,
        y + 0.28,
        3.4,
        0.3,
        size=12,
        bold=True,
        color=TEAL,
        font="Courier New",
    )
    y += 0.66

# Slide 12 — Key Files Table
sl = cslide("Key Files and Their Roles")
tbl(
    sl,
    [
        ["File", "Layer", "Role"],
        [
            "core/gridworld.py",
            "L1 Core",
            "Grid simulation, step loop, capture detection",
        ],
        ["core/agent.py", "L1 Core", "Agent state, movement, 5-action map"],
        ["observations/*.py", "L2 Plugin", "Transform env state to per-agent percepts"],
        ["rewards/*.py", "L2 Plugin", "Compute per-agent reward signals"],
        [
            "registry/observation_registry.py",
            "L2 Reg",
            "Map YAML string keys to builder classes",
        ],
        [
            "baselines/IQL/iql.py",
            "L3 Baseline",
            "Independent Q-Learning — one Q-table per agent",
        ],
        [
            "baselines/CQL/cql.py",
            "L3 Baseline",
            "Centralised Q-Learning — joint action Q-tensor",
        ],
        ["baselines/MIXED/mixed.py", "L3 Baseline", "Per-team algorithm assignment"],
        [
            "scripts/run_from_config.py",
            "L4 Script",
            "Entry point — loads configs, wires env, trains",
        ],
        [
            "scripts/evaluate.py",
            "L4 Script",
            "Greedy rollout + per-agent metrics collection",
        ],
    ],
    0.4,
    1.05,
    12.5,
    6.15,
    hfill=NAVY,
    col_w=[0.37, 0.17, 0.46],
    fs=13,
)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 5 — Class & Code Structure
# ══════════════════════════════════════════════════════════════════════════════
section(
    "Section 5: Class & Code Structure", "The key abstractions and how they connect"
)

# Slide 13 — Core Classes
sl = cslide("Core Classes — GridWorldEnv and Agent")
rect(sl, 0.4, 1.05, 6.0, 5.95, fill=RGBColor(0xFE, 0xF5, 0xF5), border=RED)
txt(sl, "GridWorldEnv(gym.Env)", 0.6, 1.1, 5.7, 0.48, size=16, bold=True, color=RED)
txt(sl, "Attributes:", 0.7, 1.65, 5.5, 0.35, size=13, bold=True, color=DGRAY)
for i, a in enumerate(
    [
        "agents              list[Agent]",
        "_obstacle_location  set of (x,y)",
        "_captured_agents    set of names",
        "reward_fn           Callable",
        "observation_builder Callable",
        "capture_threshold   int",
        "max_steps           int  |  seed  int",
    ]
):
    txt(sl, a, 0.7, 2.05 + i * 0.33, 5.5, 0.3, size=12, color=DGRAY, font="Courier New")
txt(sl, "Methods:", 0.7, 4.4, 5.5, 0.35, size=13, bold=True, color=DGRAY)
for i, m in enumerate(
    [
        "reset(seed)  →  obs_dict",
        "step(actions)  →  result_dict",
        "base_reward()  →  float",
        "_build_observations()",
    ]
):
    txt(sl, m, 0.7, 4.82 + i * 0.34, 5.5, 0.3, size=12, color=NAVY, font="Courier New")

rect(sl, 7.0, 1.05, 6.0, 5.95, fill=RGBColor(0xF5, 0xF8, 0xFF), border=BLUE)
txt(sl, "Agent", 7.2, 1.1, 5.7, 0.48, size=16, bold=True, color=BLUE)
txt(sl, "Attributes:", 7.3, 1.65, 5.5, 0.35, size=13, bold=True, color=DGRAY)
for i, a in enumerate(
    [
        "agent_name     str",
        "agent_type     'predator' | 'prey'",
        "agent_team     str",
        "_agent_location  [x, y]",
        "is_captured    bool",
        "agent_speed    int   (not yet wired)",
        "stamina        float (not yet wired)",
    ]
):
    txt(sl, a, 7.3, 2.05 + i * 0.35, 5.5, 0.3, size=12, color=DGRAY, font="Courier New")
txt(sl, "Action Map:", 7.3, 4.55, 5.5, 0.35, size=13, bold=True, color=DGRAY)
txt(
    sl,
    "0=Right  1=Up  2=Left  3=Down  4=Noop",
    7.3,
    4.95,
    5.5,
    0.32,
    size=12,
    color=NAVY,
    font="Courier New",
)
txt(
    sl,
    "All agents move simultaneously — not turn-by-turn.",
    7.3,
    5.4,
    5.5,
    0.38,
    size=12,
    italic=True,
    color=MGRAY,
)

# Slide 14 — Plugin Interfaces
sl = cslide("Plugin Interfaces — Three Abstract Base Classes")
cw = 4.05
for i, (col, title, src, note) in enumerate(
    [
        (
            GREEN,
            "ObservationBuilder",
            "class ObservationBuilder:\n  def __init__(self, **params):\n      ...\n\n  def build(\n      self,\n      env   # GridWorldEnv\n  ) -> Dict[str, dict]:\n      # Read env state\n      # Return per-agent\n      # observation dicts\n      # READ-ONLY\n      ...",
            "env.observation_builder =\n  builder_instance.build\nCalled each step.",
        ),
        (
            ORANGE,
            "RewardFunction",
            "class RewardFunction:\n  def __init__(self,\n               weight=1.0):\n      self.weight = weight\n\n  def compute(\n      self,\n      env   # GridWorldEnv\n  ) -> Dict[str, float]:\n      # Return per-agent\n      # reward contribution\n      # × self.weight\n      # READ-ONLY\n      ...",
            "Multiple functions combined\ninto one closure.\nWeights from YAML.",
        ),
        (
            PURPLE,
            "BaseAlgorithm",
            "class BaseAlgorithm:\n  def __init__(self,\n               env, config):\n      ...\n\n  def select_actions(\n      self, obs: dict\n  ) -> dict:  # required\n\n  def train(self):\n      ...   # required\n\n  def evaluate(self,\n      episodes, max_steps):\n      ...   # inherited",
            "Only env.reset() and\nenv.step() allowed.\nNo internal reads.",
        ),
    ]
):
    x = 0.4 + i * (cw + 0.2)
    rect(sl, x, 1.05, cw, 5.65, fill=RGBColor(0xF8, 0xF8, 0xF8), border=col)
    rect(sl, x, 1.05, cw, 0.48, fill=col)
    txt(sl, title, x + 0.12, 1.09, cw - 0.24, 0.4, size=14, bold=True, color=WHITE)
    code(sl, src, x, 1.53, cw, 3.9, size=11)
    rect(sl, x, 5.45, cw, 1.25, fill=RGBColor(0xF0, 0xF0, 0xF5), border=col)
    txt(sl, note, x + 0.12, 5.52, cw - 0.24, 1.1, size=12, italic=True, color=MGRAY)

# Slide 15 — Black-Box Contract
sl = cslide("The Black-Box Algorithm Contract")
txt(
    sl,
    "What the algorithm sees — and what it must never access.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
rect(sl, 0.4, 1.45, 5.95, 5.3, fill=RGBColor(0xF0, 0xFB, 0xF4), border=GREEN)
txt(sl, "✓  Allowed", 0.65, 1.5, 5.5, 0.48, size=18, bold=True, color=GREEN)
y = 2.1
for a in [
    "env.reset(seed=n)",
    "env.step(actions_dict)",
    "result['obs']",
    "result['reward']",
    "result['terminated']",
    "result['truncated']",
    "result['info']",
    "env.observation_space",
    "env.action_space",
]:
    txt(sl, "✓  " + a, 0.7, y, 5.4, 0.36, size=13, color=GREEN, font="Courier New")
    y += 0.42

rect(sl, 6.9, 1.45, 5.95, 5.3, fill=RGBColor(0xFE, 0xF5, 0xF5), border=RED)
txt(sl, "✗  Must Never Touch", 7.15, 1.5, 5.5, 0.48, size=18, bold=True, color=RED)
y = 2.1
for a in [
    "env.agents",
    "env._obstacle_location",
    "env._captured_agents",
    "agent._agent_location",
    "agent.is_captured",
    "env.reward_fn",
    "env.observation_builder",
]:
    txt(sl, "✗  " + a, 7.2, y, 5.5, 0.36, size=13, color=RED, font="Courier New")
    y += 0.52

rect(sl, 0.4, 6.82, 12.5, 0.55, fill=NAVY)
txt(
    sl,
    "The rule: any algorithm written here must work with any env that implements reset() + step(). That is what makes baselines reusable across all experimental conditions.",
    0.65,
    6.88,
    12.0,
    0.44,
    size=12,
    color=WHITE,
)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 6 — RL Theory Connection
# ══════════════════════════════════════════════════════════════════════════════
section("Section 6: Connection to RL Theory", "From Bellman equations to running code")

# Slide 16 — MDP Mapping
sl = cslide("Single-Agent MDP — Every Component Maps to Code")
txt(
    sl,
    "A Markov Decision Process is (S, A, R, P, γ). Here is exactly where each lives in this project.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
tbl(
    sl,
    [
        ["MDP", "Symbol", "In This Codebase", "Source File"],
        [
            "State Space",
            "S",
            "Output of ObservationBuilder.build(env)",
            "observations/*.py",
        ],
        [
            "Action Space",
            "A",
            "5-action discrete: {Right, Up, Left, Down, Noop}",
            "core/agent.py",
        ],
        [
            "Reward Function",
            "R(s,a)",
            "RewardFunction.compute(env) + base_reward()",
            "rewards/*.py + core/",
        ],
        [
            "Transition Model",
            "P(s'|s,a)",
            "GridWorldEnv.step() — deterministic given seed",
            "core/gridworld.py",
        ],
        [
            "Discount Factor",
            "γ",
            "config['experiment']['gamma']  (default: 0.95)",
            "configs/experiment.yaml",
        ],
        [
            "Policy",
            "π(a|s)",
            "Q-table + epsilon-greedy in select_actions()",
            "baselines/IQL/iql.py",
        ],
        [
            "Value Function",
            "Q(s,a)",
            "defaultdict — states init to zero on first visit",
            "baselines/IQL/iql.py",
        ],
    ],
    0.4,
    1.5,
    12.5,
    5.2,
    hfill=NAVY,
    col_w=[0.17, 0.09, 0.44, 0.30],
    fs=13,
)

# Slide 17 — Bellman Equation
sl = cslide("The Bellman Equation — Math to Code")
txt(
    sl,
    "The update rule that drives all tabular Q-learning in this project.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
rect(sl, 0.5, 1.42, 12.33, 1.05, fill=LGRAY, border=RGBColor(0xBB, 0xBB, 0xBB))
txt(
    sl,
    "Q(s, a)  ←  Q(s, a)  +  α  ·  [ r  +  γ · max  Q(s', a')  −  Q(s, a) ]",
    0.7,
    1.52,
    12.0,
    0.82,
    size=26,
    bold=True,
    color=NAVY,
    align=PP_ALIGN.CENTER,
)

terms = [
    (BLUE, "α  (alpha)", "Learning rate\nHow fast we update\nDefault: 0.1", 0.5),
    (GREEN, "r", "Immediate reward\nresult['reward']\n[agent_name]", 3.0),
    (TEAL, "γ  (gamma)", "Discount factor\nFuture reward weight\nDefault: 0.95", 5.5),
    (ORANGE, "max Q(s', ·)", "Best Q-value in\nnext state across\nall actions", 8.0),
    (RED, "− Q(s, a)", "Removes current\nestimate — gives\nthe TD error", 10.5),
]
y2 = 2.65
for col, label, desc, x in terms:
    rect(sl, x, y2, 2.2, 2.1, fill=col)
    txt(
        sl,
        label,
        x + 0.1,
        y2 + 0.08,
        2.0,
        0.5,
        size=15,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    txt(
        sl,
        desc,
        x + 0.1,
        y2 + 0.65,
        2.0,
        1.3,
        size=12,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

code(
    sl,
    "td_error = reward + gamma * max(q_tables[agent][s_next].values()) - q_tables[agent][s][a]\nq_tables[agent][s][a] += alpha * td_error",
    0.5,
    5.05,
    12.33,
    1.35,
    size=15,
)
txt(
    sl,
    "q_tables is a defaultdict(lambda: defaultdict(float)) — new states auto-init to zero.",
    0.7,
    6.55,
    12.0,
    0.45,
    size=13,
    italic=True,
    color=MGRAY,
)

# Slide 18 — MDP → Markov Game
sl = cslide("From MDP to Markov Game — The MARL Jump")
txt(
    sl,
    "Adding more agents changes the problem structure — not just its size.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
rect(sl, 0.4, 1.45, 5.95, 4.9, fill=RGBColor(0xF5, 0xF8, 0xFF), border=BLUE)
txt(
    sl,
    "Single-Agent MDP",
    0.6,
    1.5,
    5.5,
    0.48,
    size=17,
    bold=True,
    color=BLUE,
    align=PP_ALIGN.CENTER,
)
y = 2.1
for s in [
    "State:   S",
    "Action:  A  (one agent)",
    "Reward:  R(s, a)",
    "Trans:   P(s' | s, a)",
    "Goal:    max Σ γᵗ Rₜ",
    "",
    "Environment is stationary:",
    "rules don't change as agent learns",
]:
    txt(sl, s, 0.65, y, 5.5, 0.36, size=14, color=DGRAY, font="Courier New")
    y += 0.42

rect(sl, 6.9, 1.45, 5.95, 4.9, fill=RGBColor(0xFF, 0xF5, 0xF0), border=ORANGE)
txt(
    sl,
    "Markov Game  (MARL)",
    7.1,
    1.5,
    5.5,
    0.48,
    size=17,
    bold=True,
    color=ORANGE,
    align=PP_ALIGN.CENTER,
)
y = 2.1
for s in [
    "State:   S  (shared by all agents)",
    "Action:  A₁ × A₂ × … × Aₙ  (joint)",
    "Reward:  Rᵢ(s, a₁,…,aₙ)  per agent",
    "Trans:   P(s' | s, a₁, …, aₙ)",
    "Goal:    each agent max Σ γᵗ Rᵢₜ",
    "",
    "Non-stationarity:",
    "from agent i's view the 'environment'",
    "keeps changing as others learn",
]:
    txt(sl, s, 7.15, y, 5.6, 0.36, size=13, color=DGRAY, font="Courier New")
    y += 0.4

rect(sl, 0.4, 6.45, 12.5, 0.85, fill=NAVY)
txt(
    sl,
    "In this env: predators and prey are adversarial teams. Each team's reward depends on the other's actions. Neither side is stationary from the other's perspective.",
    0.65,
    6.52,
    12.0,
    0.72,
    size=13,
    color=WHITE,
)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 7 — MARL Algorithm Spectrum
# ══════════════════════════════════════════════════════════════════════════════
section(
    "Section 7: The MARL Algorithm Spectrum",
    "How the architecture enables different coordination strategies",
)

# Slide 19 — Spectrum
sl = cslide("The Algorithm Spectrum — Decentralised to Centralised")
txt(
    sl,
    "Each algorithm makes a different trade-off between coordination ability and scalability.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
# Spectrum bar
for i, c in enumerate([BLUE, TEAL, ORANGE, RED]):
    rect(sl, 0.4 + i * 3.23, 1.52, 3.23, 0.4, fill=c)
txt(
    sl,
    "Decentralised ←──────────────────────────────────────────────────── Centralised",
    0.42,
    1.54,
    12.46,
    0.36,
    size=13,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)

algos = [
    (
        BLUE,
        "IQL\n(Independent Q-Learning)",
        "One Q-table per agent.\nEach learns as if alone.\nIgnores teammates.",
        "✓  Scales to any n agents\n✓  Simple to implement\n✗  Non-stationarity\n✗  Coordination emergent only",
    ),
    (
        TEAL,
        "MixedTrainer\n(Per-Team)",
        "Predator team uses one algo,\nprey team uses another.\nDecision per team.",
        "✓  Team-level coordination\n✓  Flexible per-team choice\n✗  No cross-team model\n✗  Each team still non-stationary",
    ),
    (
        ORANGE,
        "CQL\n(Centralised Q-Learning)",
        "One Q-table over the joint\naction space A₁×A₂×…×Aₙ.\nFinds true joint optimum.",
        "✓  True joint optimisation\n✓  No non-stationarity\n✗  Scales as 5ⁿ actions\n✗  Impractical > 4 agents",
    ),
    (
        RED,
        "Nash Q-Learning\n(1v1 — Week 9)",
        "Game-theoretic equilibrium.\nEach agent best-responds to\nopponent's mixed strategy.",
        "✓  Game-theoretic optimality\n✓  Handles adversarial fully\n✗  1v1 only (needs LP solver)\n✗  General-sum is NP-hard",
    ),
]
x = 0.4
for col, title, desc, pros in algos:
    rect(sl, x, 2.05, 3.1, 5.0, fill=RGBColor(0xF8, 0xF8, 0xF8), border=col)
    rect(sl, x, 2.05, 3.1, 0.55, fill=col)
    txt(
        sl,
        title,
        x + 0.1,
        2.08,
        2.9,
        0.5,
        size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    txt(sl, desc, x + 0.12, 2.68, 2.86, 1.3, size=12, color=DGRAY)
    txt(sl, pros, x + 0.12, 4.1, 2.86, 1.9, size=11, color=DGRAY)
    x += 3.23

# Slide 20 — Why Architecture Enables This
sl = cslide("Why the Architecture Makes Extensibility Effortless")
txt(
    sl,
    "Adding CQL and MixedTrainer required: one file each, one registry line each, one YAML each. Zero core changes.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
examples = [
    (
        ORANGE,
        "Adding CQL to an IQL-only repo",
        "Only IQL.py in baselines/",
        "+ CQL.py in baselines/CQL/\n+ one line in algorithm_registry.py\n+ experiment_cql.yaml in configs/\n= zero changes to core, obs, rewards, or IQL",
    ),
    (
        TEAL,
        "Adding MixedTrainer",
        "Only IQL.py + CQL.py",
        "+ mixed.py in baselines/MIXED/\n+ one line in algorithm_registry.py\n+ experiment_mixed.yaml\n= zero changes anywhere else",
    ),
    (
        PURPLE,
        "Adding Nash Q-Learning (Week 9)",
        "IQL + CQL + MixedTrainer",
        "+ nashq.py in baselines/NASHQ/\n+ one line in algorithm_registry.py\n+ experiment_nashq.yaml\n= zero changes to prior baselines or core",
    ),
]
y = 1.5
for col, title, before, after in examples:
    rect(sl, 0.4, y, 12.5, 1.65, fill=RGBColor(0xF9, 0xF9, 0xF9), border=col)
    txt(sl, title, 0.65, y + 0.06, 6.0, 0.42, size=15, bold=True, color=col)
    rect(sl, 0.6, y + 0.56, 5.5, 0.95, fill=RGBColor(0xFE, 0xF5, 0xF5), border=RED)
    txt(sl, "Before: " + before, 0.75, y + 0.62, 5.2, 0.82, size=12, color=DGRAY)
    txt(
        sl,
        "→",
        6.2,
        y + 0.85,
        0.5,
        0.4,
        size=22,
        bold=True,
        color=col,
        align=PP_ALIGN.CENTER,
    )
    rect(sl, 6.8, y + 0.56, 5.95, 0.95, fill=RGBColor(0xF0, 0xFB, 0xF4), border=GREEN)
    txt(sl, "After: " + after, 6.95, y + 0.62, 5.65, 0.82, size=12, color=DGRAY)
    y += 1.82

rect(sl, 0.4, 6.82, 12.5, 0.55, fill=NAVY)
txt(
    sl,
    "The invariant: adding an algorithm never changes the environment, plugins, or other algorithms. Scientific comparability is always preserved.",
    0.65,
    6.88,
    12.0,
    0.44,
    size=12,
    color=WHITE,
)

# Slide 21 — Scaling Wall
sl = cslide("The Scaling Wall — When CQL Breaks Down")
txt(
    sl,
    "Joint action spaces grow exponentially. This is why deep MARL methods exist.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
tbl(
    sl,
    [
        ["Setup", "Agents (n)", "IQL Q-tables", "CQL Joint Actions", "CQL Practical?"],
        ["1v1", "2", "2 × |S|", "5² = 25", "Yes"],
        ["1v2", "3", "3 × |S|", "5³ = 125", "Yes"],
        ["2v2", "4", "4 × |S|", "5⁴ = 625", "Yes (barely)"],
        ["3v3", "6", "6 × |S|", "5⁶ = 15,625", "No"],
        ["5v5", "10", "10 × |S|", "5¹⁰ ≈ 10,000,000", "Never"],
    ],
    0.4,
    1.5,
    8.4,
    4.3,
    hfill=NAVY,
    col_w=[0.12, 0.15, 0.25, 0.28, 0.20],
    fs=14,
)

rect(sl, 9.1, 1.5, 3.9, 2.35, fill=LGRAY, border=NAVY)
txt(
    sl,
    "Joint space grows as:",
    9.3,
    1.56,
    3.5,
    0.45,
    size=14,
    bold=True,
    color=NAVY,
    align=PP_ALIGN.CENTER,
)
txt(
    sl, "dⁿ", 9.3, 2.05, 3.5, 0.85, size=52, bold=True, color=RED, align=PP_ALIGN.CENTER
)
txt(
    sl,
    "d = 5 actions  •  n = agents",
    9.3,
    2.98,
    3.5,
    0.45,
    size=13,
    color=DGRAY,
    align=PP_ALIGN.CENTER,
)

rect(sl, 9.1, 4.0, 3.9, 1.75, fill=RGBColor(0xFF, 0xF5, 0xF0), border=ORANGE)
txt(
    sl,
    "Why IQL survives:\nEach agent has its OWN\nQ-table — no joint space.\nO(n · |S| · d)  not  O(|S| · dⁿ)",
    9.3,
    4.1,
    3.6,
    1.55,
    size=13,
    color=DGRAY,
)

rect(sl, 0.4, 6.0, 12.5, 1.3, fill=RGBColor(0xFF, 0xF8, 0xF0), border=ORANGE)
txt(
    sl,
    "This project: optimised for ≤ 4 agents on grids ≤ 10×10.\nFor larger games, the natural next step is CTDE (Centralised Training, Decentralised Execution) methods\nsuch as QMIX or MADDPG — deep RL baselines planned for future work.",
    0.65,
    6.08,
    12.0,
    1.18,
    size=13,
    color=DGRAY,
)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 8 — Contribution Policy
# ══════════════════════════════════════════════════════════════════════════════
section("Section 8: Contributing", "Rules, pathways, and checklist")

# Slide 22 — Golden Rule
sl = cslide("The Golden Rule — Immutable Core")
txt(
    sl,
    "One rule protects the scientific integrity of every experiment run on this codebase.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
rect(sl, 0.9, 1.42, 11.53, 1.45, fill=NAVY, border=ORANGE)
txt(
    sl,
    "core/gridworld.py  and  core/agent.py  are NEVER modified by contributors.",
    1.1,
    1.5,
    11.1,
    0.52,
    size=22,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)
txt(
    sl,
    "All variation is expressed through config files, observation plugins, reward plugins, and algorithm plugins.",
    1.1,
    2.12,
    11.1,
    0.65,
    size=16,
    color=AMBER,
    align=PP_ALIGN.CENTER,
)
txt(sl, "Why?", 0.5, 3.2, 12.0, 0.45, size=18, bold=True, color=NAVY)
bullets(
    sl,
    [
        "Any physics change invalidates comparisons across ALL prior experiments — you lose the baseline",
        "The invariant is what makes this a fair test bed: change the rules, lose the comparison",
        "Every variation you need can be expressed through the three plugin layers",
        "If it cannot, that is a discussion for maintainers — not a unilateral edit",
        "Enforced by: code review + CONTRIBUTING.md + ADR-001",
    ],
    0.5,
    3.72,
    12.33,
    3.0,
    size=16,
)

# Slide 23 — Three Pathways
sl = cslide("The Three Contribution Pathways")
txt(
    sl,
    "Every contribution fits one of these patterns. Following the pathway = guaranteed clean merge.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
pathways = [
    (
        GREEN,
        "Pathway 1: New Observation Builder",
        "1. src/observations/my_obs.py  — subclass ObservationBuilder, implement build(env)\n"
        "2. Add one line to  observation_registry.py\n"
        "3. Add YAML key to  configs/observations.yaml\n"
        "4. Write tests in  tests/test_observations.py",
    ),
    (
        ORANGE,
        "Pathway 2: New Reward Function",
        "1. src/rewards/my_reward.py  — subclass RewardFunction, implement compute(env)\n"
        "2. Add one line to  reward_registry.py\n"
        "3. Add YAML block to  configs/rewards.yaml\n"
        "4. Write tests in  tests/test_rewards.py",
    ),
    (
        PURPLE,
        "Pathway 3: New Learning Algorithm",
        "1. src/baselines/MYALGO/myalgo.py  — subclass BaseAlgorithm, implement select_actions() + train()\n"
        "2. Add self-registration guard at module level\n"
        "3. Import in  baselines/__init__.py\n"
        "4. Add  configs/experiment_myalgo.yaml\n"
        "5. Write tests in  tests/test_baselines_myalgo.py",
    ),
]
y = 1.5
for col, title, steps in pathways:
    h = 1.5 if col != PURPLE else 1.85
    rect(sl, 0.4, y, 12.5, h, fill=RGBColor(0xF8, 0xF8, 0xF8), border=col)
    rect(sl, 0.4, y, 12.5, 0.46, fill=col)
    txt(sl, title, 0.65, y + 0.06, 12.0, 0.38, size=15, bold=True, color=WHITE)
    code(sl, steps, 0.4, y + 0.47, 12.5, h - 0.48, size=11)
    y += h + 0.12


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 9 — End-to-End Example
# ══════════════════════════════════════════════════════════════════════════════
section(
    "Section 9: End-to-End Example",
    "One command, traceable from YAML to trained policy",
)

# Slide 24 — The Scenario
sl = cslide("The Scenario — Train Predators to Capture Prey")
txt(
    sl,
    "Goal: 2 predators learn to capture 2 prey on a 6×6 grid using IQL + distance shaping.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=15,
    italic=True,
    color=MGRAY,
)

code(
    sl,
    "# env.yaml\ngrid_size: 6\nmax_steps: 200\ncapture_threshold: 1\nperc_num_obstacle: 15.0\nrender_mode: null\nseed: 42",
    0.4,
    1.4,
    3.8,
    2.4,
    size=12,
)
code(
    sl,
    "# agents.yaml\nagents:\n  - name: pred_0\n    type: predator\n    team: predators\n  - name: pred_1\n    type: predator\n    team: predators\n  - name: prey_0\n    type: prey\n    team: prey\n  - name: prey_1\n    type: prey\n    team: prey",
    4.4,
    1.4,
    3.8,
    4.0,
    size=12,
)
code(
    sl,
    "# observations.yaml\nobservation:\n  type: local_radius\n  params:\n    radius: 3\n\n# rewards.yaml\nrewards:\n  - name: base_reward\n    weight: 1.0\n  - name: predator_distance\n    weight: 0.5\n\n# experiment.yaml\nexperiment:\n  algorithm:\n    name: iql\n  episodes: 1000\n  alpha: 0.1\n  gamma: 0.95\n  epsilon: 1.0\n  epsilon_decay: 0.995\n  min_epsilon: 0.01",
    8.4,
    1.4,
    4.55,
    5.4,
    size=11,
)

code(
    sl,
    "# ONE COMMAND — no training code written\npython -m multi_agent_package.scripts.run_from_config --config configs/",
    0.4,
    3.95,
    7.8,
    0.9,
    size=14,
)
txt(
    sl,
    "The only code you write is config. The system assembles everything.",
    0.6,
    5.05,
    7.6,
    0.42,
    size=13,
    italic=True,
    color=MGRAY,
)

# Slide 25 — From Config to Training
sl = cslide("One Command — Full Training Run")
txt(
    sl,
    "What happens when you hit Enter:",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
rect(sl, 0.4, 1.38, 12.5, 0.72, fill=NAVY)
txt(
    sl,
    "python -m multi_agent_package.scripts.run_from_config  --config configs/",
    0.65,
    1.45,
    12.0,
    0.58,
    size=18,
    bold=True,
    color=AMBER,
    font="Courier New",
)

e2e = [
    (BLUE, "① Load", "5 YAML files\n→ merged config"),
    (TEAL, "② Agents", "4 Agent objects\npred_0…prey_1"),
    (GREEN, "③ Env", "GridWorldEnv\nsize=6, seed=42"),
    (ORANGE, "④ Obs", "LocalRadius(r=3)\n→ env.obs_builder"),
    (PURPLE, "⑤ Reward", "Base + Distance\n→ env.reward_fn"),
    (RED, "⑥ Train", "IQL.train()\n1000 episodes"),
]
bw2 = 12.5 / 6
for i, (col, step, desc) in enumerate(e2e):
    x = 0.4 + i * bw2
    rect(sl, x, 2.3, bw2 - 0.08, 0.5, fill=col)
    txt(
        sl,
        step,
        x + 0.08,
        2.34,
        bw2 - 0.24,
        0.42,
        size=14,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    rect(sl, x, 2.8, bw2 - 0.08, 1.05, fill=LGRAY, border=col)
    txt(
        sl,
        desc,
        x + 0.08,
        2.87,
        bw2 - 0.24,
        0.92,
        size=12,
        color=DGRAY,
        align=PP_ALIGN.CENTER,
    )

rect(sl, 0.4, 4.1, 12.5, 3.05, fill=RGBColor(0x1E, 0x1E, 0x1E))
txt(sl, "Console output:", 0.6, 4.16, 5.0, 0.32, size=12, bold=True, color=MGRAY)
txt(
    sl,
    (
        "Episode  100/1000  |  avg_reward: -12.4  |  captures: 0.12  |  epsilon: 0.607\n"
        "Episode  200/1000  |  avg_reward:  -8.1  |  captures: 0.31  |  epsilon: 0.368\n"
        "Episode  500/1000  |  avg_reward:   2.7  |  captures: 0.64  |  epsilon: 0.082\n"
        "Episode 1000/1000  |  avg_reward:  14.2  |  captures: 0.88  |  epsilon: 0.010\n"
        "Training complete. Saved: checkpoints/iql_pred_0.pkl  (+ 3 others)"
    ),
    0.6,
    4.55,
    12.0,
    2.45,
    size=13,
    color=RGBColor(0x00, 0xFF, 0x41),
    font="Courier New",
)

# Slide 26 — Observing the Result
sl = cslide("What Success Looks Like")
txt(
    sl,
    "Same seed + same YAMLs reproduces this result exactly on any machine.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)
tbl(
    sl,
    [
        ["Metric", "Trained Policy", "Random Policy"],
        ["Capture rate", "88%", "12%"],
        ["Avg episode length", "67 steps", "198 steps"],
        ["Avg total reward (pred)", "+14.2", "-8.6"],
        ["Avg total reward (prey)", "-31.0", "-4.1"],
        ["Unique states visited", "1,847", "312"],
    ],
    0.4,
    1.5,
    5.7,
    3.9,
    hfill=TEAL,
    col_w=[0.48, 0.27, 0.25],
    fs=14,
)
code(
    sl,
    "python -m multi_agent_package.scripts.evaluate \\\n  --checkpoint checkpoints/ --episodes 200",
    0.4,
    5.6,
    5.7,
    0.95,
    size=13,
)

rect(sl, 6.4, 1.5, 6.7, 4.8, fill=RGBColor(0xF8, 0xF8, 0xFF), border=BLUE)
txt(sl, "How to read the results", 6.6, 1.56, 6.3, 0.45, size=14, bold=True, color=NAVY)
insights = [
    (
        "Reward curve rising",
        "Predators learned to move toward prey.\nDistance shaping drove this behaviour.",
    ),
    (
        "Episode length falling",
        "Captures happen earlier as policy improves.\nRandom agents always hit max_steps.",
    ),
    (
        "States visited growing",
        "Q-table coverage expands during exploration\nthen stabilises when epsilon is low.",
    ),
    (
        "Prey reward drops too",
        "Prey Q-tables also learn evasion —\nbut predators win this game.",
    ),
]
y = 2.1
for title, body in insights:
    txt(sl, "▶  " + title, 6.6, y, 6.3, 0.36, size=14, bold=True, color=BLUE)
    txt(sl, body, 6.6, y + 0.38, 6.3, 0.62, size=13, color=DGRAY)
    y += 1.1


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Status + Roadmap
# ══════════════════════════════════════════════════════════════════════════════
sl = cslide("Where We Are — and Where We're Going")
txt(
    sl,
    "Ten weeks, four students, one JOSS paper.",
    0.5,
    0.92,
    12.33,
    0.42,
    size=14,
    italic=True,
    color=MGRAY,
)

rect(sl, 0.4, 1.45, 5.85, 5.35, fill=RGBColor(0xF0, 0xFB, 0xF4), border=GREEN)
txt(
    sl,
    "Current State  (start of Week 1)",
    0.6,
    1.5,
    5.5,
    0.45,
    size=15,
    bold=True,
    color=GREEN,
)
y = 2.05
for s in [
    "Discrete gridworld (up to 10×10)",
    "5 observation builders",
    "3 reward functions",
    "IQL + CQL + MixedTrainer (tabular)",
    "189 pytest tests across 9 files",
    "Full SDD wiki (38 documents)",
    "Interactive notebooks (IQL + CQL)",
    "CI pipeline (in progress)",
]:
    txt(sl, "✓  " + s, 0.65, y, 5.3, 0.38, size=14, color=DGRAY)
    y += 0.52

rect(sl, 6.6, 1.45, 6.35, 5.35, fill=RGBColor(0xF5, 0xF8, 0xFF), border=BLUE)
txt(sl, "10-Week Roadmap", 6.8, 1.5, 5.9, 0.45, size=15, bold=True, color=BLUE)
roadmap = [
    ("Week 1", "Bug fixes • Gymnasium wrapper • CI passing"),
    ("Weeks 2–4", "IDQN — neural Q-learning per agent"),
    ("Weeks 5–6", "IPPO — actor-critic, policy gradient"),
    ("Weeks 7–9", "Nash Q-Learning — 1v1 minimax"),
    ("Week 10", "JOSS paper + final documentation"),
]
y = 2.05
for week, item in roadmap:
    rect(sl, 6.8, y, 5.9, 0.78, fill=LGRAY, border=RGBColor(0xBB, 0xBB, 0xBB))
    txt(
        sl,
        week,
        6.95,
        y + 0.1,
        1.2,
        0.38,
        size=12,
        bold=True,
        color=NAVY,
        font="Courier New",
    )
    txt(sl, item, 8.2, y + 0.1, 4.35, 0.62, size=13, color=DGRAY)
    y += 0.9

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Resources (dark closing slide)
# ══════════════════════════════════════════════════════════════════════════════
sl = S()
bg(sl, NAVY)
rect(sl, 0, 0, 13.33, 0.32, fill=ORANGE)
txt(
    sl,
    "Resources & Next Steps",
    0.5,
    0.55,
    12.33,
    0.85,
    size=34,
    bold=True,
    color=WHITE,
    align=PP_ALIGN.CENTER,
)
rect(sl, 1.0, 1.55, 11.33, 0.06, fill=ORANGE)

cards = [
    (
        BLUE,
        "Start Here",
        [
            "wiki/index.md  — documentation entry point",
            "notebooks/01_iql_demo.ipynb  — interactive IQL walk-through",
            "QUICKSTART.md  — install + first run in 5 minutes",
        ],
    ),
    (
        TEAL,
        "Key References",
        [
            "Sutton & Barto (2018)  — Chapters 6, 13",
            "Hu & Wellman (2003)  — Nash Q-Learning",
            "Littman (2001)  — Friend-or-Foe Q-Learning",
        ],
    ),
    (
        ORANGE,
        "This Week",
        [
            "wiki/weekly/week-01.md  — full task list with owners",
            "wiki/reviews/audit-2026-06-07.md  — known bugs",
            "CONTRIBUTING.md  — pull request checklist",
        ],
    ),
]
x = 0.5
for col, title, items in cards:
    rect(sl, x, 1.85, 4.0, 5.2, fill=RGBColor(0x0D, 0x22, 0x38), border=col)
    rect(sl, x, 1.85, 4.0, 0.5, fill=col)
    txt(sl, title, x + 0.15, 1.88, 3.7, 0.42, size=15, bold=True, color=WHITE)
    y = 2.5
    for item in items:
        txt(
            sl,
            "•  " + item,
            x + 0.15,
            y,
            3.7,
            0.62,
            size=12,
            color=RGBColor(0xCC, 0xDD, 0xEE),
        )
        y += 0.75
    x += 4.4

txt(
    sl,
    "Build understanding from the ground up.",
    0.5,
    7.1,
    12.33,
    0.35,
    size=16,
    italic=True,
    color=AMBER,
    align=PP_ALIGN.CENTER,
)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = Path(__file__).parent.parent / "slides" / "presentation.pptx"
out.parent.mkdir(exist_ok=True)
prs.save(str(out))
print(f"Saved {len(prs.slides)} slides -> {out}")
